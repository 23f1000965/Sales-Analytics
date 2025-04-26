import streamlit as st
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
import plotly.express as px
from datetime import datetime
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Set Page Config
st.set_page_config(page_title="Sales Insights Dashboard", layout="wide")

# Setup Theme
sns.set_theme(style="whitegrid", palette="pastel")

# Temporary folder for saving graphs
TEMP_FOLDER = "tmp_images"
os.makedirs(TEMP_FOLDER, exist_ok=True)

# Load Data
@st.cache_data
def load_data(uploaded_file):
    if uploaded_file is not None:
        df = pd.read_excel(uploaded_file)
        df.columns = df.columns.str.strip()
        df['order_date'] = pd.to_datetime(df['Order Date'], errors='coerce')
        df['total_revenue'] = pd.to_numeric(df['Invoice Amount'], errors='coerce')
        return df
    else:
        st.warning("Upload a file to continue.")
        return None

uploaded_file = st.file_uploader("Upload your Excel file", type=["xlsx"])
df = load_data(uploaded_file)

if df is None:
    st.stop()

# Sidebar Filters
st.sidebar.header("Filters")
date_range = st.sidebar.date_input("Select Date Range", [df['order_date'].min(), df['order_date'].max()])
products = st.sidebar.multiselect("Select Products", options=df['Item Description'].unique())
states = st.sidebar.multiselect("Select States", options=df['Ship To State'].unique())


filtered_df = df.copy()
if date_range:
    filtered_df = filtered_df[(filtered_df['order_date'] >= pd.to_datetime(date_range[0])) & (filtered_df['order_date'] <= pd.to_datetime(date_range[1]))]
if products:
    filtered_df = filtered_df[filtered_df['Item Description'].isin(products)]
if states:
    filtered_df = filtered_df[filtered_df['Ship To State'].isin(states)]


# Overview Tab
st.title("📊 Sales Insights Dashboard")
tab1, tab2, tab3, tab4 = st.tabs(["Overview", "Products", "Shipping", "Returns"])

with tab1:
    st.subheader("Daily Sales Trend")
    daily_sales = filtered_df.groupby('order_date')['total_revenue'].sum().reset_index()
    fig, ax = plt.subplots(figsize=(10, 5))
    sns.lineplot(data=daily_sales, x='order_date', y='total_revenue', ax=ax)
    ax.set_ylabel("Revenue (₹)")
    plt.xticks(rotation=45)
    fig.tight_layout()
    st.pyplot(fig)
    fig.savefig(f"{TEMP_FOLDER}/daily_sales_trend.png")

    st.subheader("State-wise Sales")
    state_sales = df.groupby('Ship To State').agg({
        'Invoice Amount': 'sum',
        'Order Id': 'nunique',
        'Quantity': 'sum'
    }).sort_values('Invoice Amount', ascending=False).reset_index()

    top_state = state_sales.iloc[0]
    st.markdown(f"### 🥇 Top State: **{top_state['Ship To State']}**")
    col4, col5, col6 = st.columns(3)
    col4.metric("Revenue", f"₹{top_state['Invoice Amount']:,.2f}")
    col5.metric("Orders", int(top_state['Order Id']))
    col6.metric("Units Sold", int(top_state['Quantity']))

    fig2 = px.bar(state_sales.head(10), x='Invoice Amount', y='Ship To State', orientation='h', color='Ship To State', title="Top 10 States by Sales")
    st.plotly_chart(fig2, use_container_width=True)
    fig2.write_image(f"{TEMP_FOLDER}/top_states_sales.png")

    top_product_name = filtered_df['Item Description'].mode()[0] if not filtered_df.empty else "N/A"
    st.markdown("#### Insights")
    st.info(f"🏆 Highest sales were recorded in **{top_state['Ship To State']}**, driven mainly by **{top_product_name}**.")

with tab2:
    st.subheader("Top Products by Revenue")
    top_products = df.groupby('Item Description')['Invoice Amount'].sum().nlargest(10).reset_index()
    st.dataframe(top_products)

    top_products_chart = filtered_df.groupby('Item Description')['total_revenue'].sum().sort_values(ascending=False).head(10)
    st.bar_chart(top_products_chart)
    
    fig3, ax3 = plt.subplots()
    top_products_chart.plot(kind='barh', ax=ax3)
    plt.title("Top Products by Revenue")
    plt.xlabel("Revenue (₹)")
    fig3.tight_layout()
    fig3.savefig(f"{TEMP_FOLDER}/top_products_revenue.png")

with tab3:
    st.subheader("Shipping Amount Trend")
    st.metric("Total Shipping Amount", f"₹{df['Shipping Amount'].sum():,.2f}")
    if 'Shipping Amount' in filtered_df.columns:
        shipping_trend = filtered_df.groupby('order_date')['Shipping Amount'].sum()
        st.line_chart(shipping_trend)

        fig4, ax4 = plt.subplots()
        shipping_trend.plot(ax=ax4)
        plt.title("Shipping Amount Trend")
        plt.ylabel("Shipping Amount (₹)")
        plt.xticks(rotation=45)
        fig4.tight_layout()
        fig4.savefig(f"{TEMP_FOLDER}/shipping_trend.png")

with tab4:
    st.subheader("Top Returned Products")
    filtered_df['Invoice Date'] = pd.to_datetime(filtered_df['Invoice Date'], errors='coerce')
    filtered_df['Invoice Day'] = filtered_df['Invoice Date'].dt.date

    returned_df = filtered_df[filtered_df['Quantity'] < 0]
    returned_products = returned_df['Item Description'].value_counts().reset_index().head(10)
    returned_products.columns = ['Item Description', 'Return Count']

    fig5, ax5 = plt.subplots()
    sns.barplot(data=returned_products, x='Return Count', y='Item Description', ax=ax5)
    plt.title("Top Returned Products")
    fig5.tight_layout()
    st.pyplot(fig5)
    fig5.savefig(f"{TEMP_FOLDER}/top_returned_products.png")

    st.subheader("Returns Over Time")
    returns_daily = filtered_df[filtered_df['Credit Note No'].notnull()].groupby('Invoice Day')['Invoice Amount'].sum().reset_index()
    fig6, ax6 = plt.subplots()
    sns.lineplot(data=returns_daily, x='Invoice Day', y='Invoice Amount', marker='o', color='red', ax=ax6)
    plt.title("Returns Over Time")
    plt.xticks(rotation=45)
    fig6.tight_layout()
    st.pyplot(fig6)
    fig6.savefig(f"{TEMP_FOLDER}/returns_over_time.png")

# -------------------------
# 📥 Download PPT

def generate_ppt():
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sales Insights Dashboard Feburary 2025"
    
    # Add Metrics Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Summary Metrics"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = f"Total Revenue: ₹{filtered_df['total_revenue'].sum():,.2f}"
    tf.add_paragraph().text = f"Total Orders: {filtered_df['Order Id'].nunique()}"
    tf.add_paragraph().text = f"Total Units Sold: {filtered_df['Quantity'].sum()}"
    tf.add_paragraph().text = f"Total Shipping Amount: ₹{df['Shipping Amount'].sum():,.2f}"

    tf.add_paragraph().text = f"Top State: {top_state['Ship To State']}"
    tf.add_paragraph().text = f"Top Product: {top_product_name}"

    # Insert Graphs Slide by Slide
    for graph_name in [
        "daily_sales_trend.png",
        "top_states_sales.png",
        "top_products_revenue.png",
        "shipping_trend.png",
        "top_returned_products.png",
        "returns_over_time.png",
    ]:
        img_path = os.path.join(TEMP_FOLDER, graph_name)
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = graph_name.replace("_", " ").replace(".png", "").title()
            slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), Inches(8), Inches(4.5))

    # Save PPT
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

import io
ppt_bytes = generate_ppt()

st.sidebar.download_button(
    label="📥 Download Sales Report PPT",
    data=ppt_bytes,
    file_name=f"sales_report_{datetime.now().strftime('%Y%m%d')}.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
)

